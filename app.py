# ================= IMPORTS =================
from flask import Flask, render_template, request, redirect, session
import mysql.connector
import requests
from bs4 import BeautifulSoup
from urllib.parse import urljoin
import os

# NEW: file processing
from werkzeug.utils import secure_filename
import PyPDF2
import docx
from pptx import Presentation

# ML Models
from transformers import pipeline
from transformers import BlipProcessor, BlipForConditionalGeneration
from PIL import Image
import whisper
import yt_dlp

app = Flask(__name__)
app.secret_key = "secretkey"

UPLOAD_FOLDER = "uploads"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER

# ================= DATABASE =================
db = mysql.connector.connect(
    host="localhost",
    user="root",
    password="",
    database="multimodal_db",
    charset="utf8"
)
cursor = db.cursor()

# ================= MODELS =================
summarizer = pipeline("text2text-generation", model="google/flan-t5-base")

processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-base")
blip_model = BlipForConditionalGeneration.from_pretrained("Salesforce/blip-image-captioning-base")

whisper_model = whisper.load_model("base")

# ================= ROUTES =================

@app.route("/")
def index():
    return render_template("index.html")


@app.route("/register", methods=["GET", "POST"])
def register():
    if request.method == "POST":
        data = (
            request.form["name"],
            request.form["email"],
            request.form["mobile"],
            request.form["username"],
            request.form["password"]
        )

        cursor.execute(
            "INSERT INTO users (name,email,mobile,username,password) VALUES (%s,%s,%s,%s,%s)",
            data
        )
        db.commit()
        return redirect("/login")

    return render_template("register.html")


@app.route("/login", methods=["GET", "POST"])
def login():
    if request.method == "POST":
        username = request.form["username"]
        password = request.form["password"]

        cursor.execute(
            "SELECT * FROM users WHERE username=%s AND password=%s",
            (username, password)
        )
        user = cursor.fetchone()

        if user:
            session["user"] = username
            return redirect("/dashboard")

        return "Invalid Login"

    return render_template("login.html")


@app.route("/dashboard", methods=["GET", "POST"])
def dashboard():
    if "user" not in session:
        return redirect("/login")

    result = ""

    if request.method == "POST":

        url = request.form.get("url", "").strip()
        file = request.files.get("file")

        text, images, videos = "", [], []

        # ================= FILE UPLOAD =================
        if file and file.filename != "":
            filename = secure_filename(file.filename)
            path = os.path.join(app.config["UPLOAD_FOLDER"], filename)
            file.save(path)

            print("Processing file:", filename)

            text = extract_file_text(path)

        # ================= URL =================
        elif url:
            # Normalize
            url = url.strip()

            # Detect YouTube (ALL formats)
            if any(x in url for x in ["youtube.com", "youtu.be"]):
                videos = [url]

            # Direct video links
            elif any(url.lower().endswith(ext) for ext in [".mp4", ".webm", ".ogg"]):
                videos = [url]

            else:
                text, images, videos = extract_content(url)

        # ================= PROCESS =================
        ST = summarize_text(text)
        SI = caption_images(images)
        SV = process_video(videos)

        result = final_summary(ST, SI, SV)

    return render_template("dashboard.html", result=result)


@app.route("/logout")
def logout():
    session.pop("user", None)
    return redirect("/")


# ================= FILE TEXT EXTRACTION =================

def extract_file_text(path):
    text = ""

    try:
        if path.endswith(".pdf"):
            with open(path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text += page.extract_text() or ""

        elif path.endswith(".docx"):
            doc = docx.Document(path)
            for para in doc.paragraphs:
                text += para.text + " "

        elif path.endswith(".pptx"):
            prs = Presentation(path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + " "

    except Exception as e:
        print("File read error:", e)

    return text


# ================= EXISTING FUNCTIONS (UNCHANGED) =================
# KEEP ALL YOUR ORIGINAL FUNCTIONS BELOW (NO CHANGE)

def extract_content(url):
    try:
        headers = {"User-Agent": "Mozilla/5.0"}
        response = requests.get(url, headers=headers, timeout=10)
        soup = BeautifulSoup(response.text, "html.parser")

        # ================= TEXT =================
        text = " ".join([p.get_text() for p in soup.find_all("p")])

        # ================= IMAGES =================
        images = []
        for img in soup.find_all("img"):
            src = img.get("src") or img.get("data-src")
            if src:
                images.append(urljoin(url, src))

        # ================= VIDEOS =================
        videos = set()

        # 1. <video> tags
        for v in soup.find_all("video"):
            src = v.get("src") or v.get("data-src")
            if src:
                videos.add(urljoin(url, src))

        # 2. <source> tags
        for s in soup.find_all("source"):
            src = s.get("src")
            if src and any(ext in src for ext in [".mp4", ".webm", ".m3u8"]):
                videos.add(urljoin(url, src))

        # 3. iframes (ANY video platform)
        for iframe in soup.find_all("iframe"):
            src = iframe.get("src")
            if src and not "google.com" in src:
                videos.add(src)

        # 4. OpenGraph meta tags (VERY IMPORTANT)
        for meta in soup.find_all("meta"):
            prop = meta.get("property", "")
            content = meta.get("content", "")

            if "video" in prop and content:
                videos.add(content)

        # 5. RAW links in page (hidden video URLs)
        for link in soup.find_all("a"):
            href = link.get("href")
            if href and any(ext in href for ext in [".mp4", ".webm", ".m3u8"]):
                videos.add(urljoin(url, href))

        return text, images, list(videos)

    except Exception as e:
        print("Error:", e)
        return "", [], []

def summarize_text(text):
    if not text.strip():
        return ""

    result = summarizer(
        "Write a very detailed explanation in at least 20 lines:\n" + text[:2000],
        max_length=400,
        min_length=200,
        do_sample=True,
        temperature=0.7,
        repetition_penalty=2.0
    )

    return result[0]["generated_text"]


def caption_images(image_urls):
    print("Detected images:", image_urls)  # DEBUG

    captions = []

    for url in image_urls[:3]:
        try:
            img = Image.open(requests.get(url, stream=True).raw).convert("RGB")
            inputs = processor(img, return_tensors="pt")
            out = blip_model.generate(**inputs)
            caption = processor.decode(out[0], skip_special_tokens=True)
            captions.append(caption)
        except Exception as e:
            print("Image error:", e)

    return " ".join(captions)

def process_video(video_urls):
    if not video_urls:
        return ""

    results = []

    for video_url in video_urls[:2]:
        try:
            print("Processing:", video_url)

            if any(ext in video_url.lower() for ext in [".mp4", ".webm", ".ogg"]):
                text = transcribe_video_file(video_url)

                if text.strip():
                    results.append(text)
                else:
                    results.append("Video has no speech (only visuals or music).")

            else:
                results.append(f"Detected video (cannot process): {video_url}")

        except Exception as e:
            print("Video error:", e)

    return " ".join(results)

def transcribe_video_file(video_url):
    try:
        file = "temp_video.mp4"

        print("Downloading video...")

        r = requests.get(video_url, stream=True, timeout=10)
        with open(file, "wb") as f:
            for chunk in r.iter_content(1024):
                f.write(chunk)

        print("Transcribing...")

        result = whisper_model.transcribe(file)

        os.remove(file)

        print("Transcription result:", result["text"])

        return result["text"]

    except Exception as e:
        print("Transcription error:", e)
        return ""

    
def final_summary(ST, SI, SV):
    content = " ".join([ST, SI, SV])

    if not content.strip():
        return "No content available"

    result = summarizer(
        "Write a very detailed explanation in more than 20 lines:\n" + content,
        max_length=500,
        min_length=250,
        do_sample=True,
        temperature=0.7,
        repetition_penalty=2.0
    )

    return result[0]["generated_text"]




# ================= ADMIN LOGIN =================
@app.route("/admin", methods=["GET", "POST"])
def admin():
    if request.method == "POST":
        username = request.form["username"]
        password = request.form["password"]

        if username == "admin" and password == "admin":
            session["admin"] = True
            return redirect("/admin/dashboard")

        return "Invalid Admin Login"

    return render_template("admin_login.html")


# ================= ADMIN DASHBOARD =================
@app.route("/admin/dashboard")
def admin_dashboard():
    if "admin" not in session:
        return redirect("/admin")

    # Total users
    cursor.execute("SELECT COUNT(*) FROM users")
    total_users = cursor.fetchone()[0]

    # Fetch users (without password)
    cursor.execute("SELECT id, name, email, mobile, username FROM users")
    users = cursor.fetchall()

    return render_template(
        "admin_dashboard.html",
        total_users=total_users,
        users=users
    )


# ================= ADMIN LOGOUT =================
@app.route("/admin/logout")
def admin_logout():
    session.pop("admin", None)
    return redirect("/")

# ================= RUN =================
if __name__ == "__main__":
    app.run(debug=False)
